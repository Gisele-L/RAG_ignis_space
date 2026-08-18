from pathlib import Path
import csv
import json
import re

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader


# ============================================================
# CONFIGURAÇÃO
# ============================================================

BASE_DIR = Path(__file__).resolve().parent.parent

DOCS_DIR = BASE_DIR / "docs"


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".csv",
    ".json",
    ".md",
    ".html",
    ".pptx",
}


# ============================================================
# PADRÕES SUSPEITOS DE PROMPT INJECTION
# ============================================================

PROMPT_INJECTION_PATTERNS = {
    "ignore_previous_instructions": (
        r"\bignore\s+(all\s+)?previous\s+instructions\b"
    ),

    "ignore_prior_instructions": (
        r"\bignore\s+(all\s+)?prior\s+instructions\b"
    ),

    "ignore_system_prompt": (
        r"\bignore\s+(the\s+)?system\s+prompt\b"
    ),

    "system_prompt_reference": (
        r"\bsystem\s+prompt\b"
    ),

    "developer_message_reference": (
        r"\bdeveloper\s+(message|instructions?)\b"
    ),

    "assistant_role_marker": (
        r"(?im)^\s*assistant\s*:"
    ),

    "system_role_marker": (
        r"(?im)^\s*system\s*:"
    ),

    "developer_role_marker": (
        r"(?im)^\s*developer\s*:"
    ),

    "you_are_chatgpt": (
        r"\byou\s+are\s+chatgpt\b"
    ),

    "do_not_follow": (
        r"\bdo\s+not\s+follow\b"
    ),

    "disregard_instructions": (
        r"\bdisregard\s+.*instructions?\b"
    ),

    "override_instructions": (
        r"\boverride\s+.*instructions?\b"
    ),
}


# ============================================================
# SEGURANÇA
# ============================================================

def detect_prompt_injection(content: str) -> dict:
    """
    Procura padrões potencialmente associados a prompt injection.

    A função não bloqueia o documento.

    Ela apenas retorna informações que serão adicionadas
    aos metadados do Document.
    """

    matches = []

    for pattern_name, pattern in PROMPT_INJECTION_PATTERNS.items():

        if re.search(
            pattern,
            content,
            flags=re.IGNORECASE
        ):

            matches.append(
                pattern_name
            )

    return {
        "prompt_injection_risk": bool(matches),
        "prompt_injection_matches": matches,
    }


# ============================================================
# FUNÇÕES AUXILIARES
# ============================================================

def create_document(
    content: str,
    source: Path,
    file_type: str,
    **metadata
):
    """
    Cria um objeto Document padronizado para o RAG.
    """

    security_metadata = detect_prompt_injection(
        content
    )

    return Document(
        page_content=content,
        metadata={
            "source": str(
                source.relative_to(BASE_DIR)
            ),
            "file_name": source.name,
            "file_type": file_type,
            **security_metadata,
            **metadata,
        },
    )


# ============================================================
# PDF
# ============================================================

def load_pdf(file_path: Path):
    """
    Carrega um PDF preservando a separação por página.
    """

    loader = PyPDFLoader(
        str(file_path)
    )

    documents = loader.load()

    for document in documents:

        security_metadata = detect_prompt_injection(
            document.page_content
        )

        document.metadata.update(
            {
                "source": str(
                    file_path.relative_to(BASE_DIR)
                ),
                "file_name": file_path.name,
                "file_type": "pdf",
                **security_metadata,
            }
        )

    return documents


# ============================================================
# DOCX
# ============================================================

def load_docx(file_path: Path):
    """
    Carrega texto de documentos DOCX.
    """

    doc = DocxDocument(
        file_path
    )

    paragraphs = []

    for paragraph in doc.paragraphs:

        text = paragraph.text.strip()

        if text:
            paragraphs.append(
                text
            )


    # Também captura conteúdo básico das tabelas
    for table in doc.tables:

        for row in table.rows:

            cells = [
                cell.text.strip()
                for cell in row.cells
            ]

            if any(cells):

                paragraphs.append(
                    " | ".join(cells)
                )


    content = "\n".join(
        paragraphs
    )


    if not content.strip():
        return []


    return [
        create_document(
            content,
            file_path,
            "docx",
        )
    ]


# ============================================================
# XLSX
# ============================================================

def load_xlsx(file_path: Path):
    """
    Converte planilhas XLSX em texto estruturado.
    """

    workbook = load_workbook(
        filename=file_path,
        data_only=True,
        read_only=True,
    )


    documents = []


    for worksheet in workbook.worksheets:

        rows = []


        for row in worksheet.iter_rows(
            values_only=True
        ):

            values = []


            for value in row:

                if value is None:
                    values.append("")

                else:
                    values.append(
                        str(value)
                    )


            if any(
                value.strip()
                for value in values
            ):

                rows.append(
                    " | ".join(values)
                )


        if not rows:
            continue


        content = (
            f"Planilha: {worksheet.title}\n\n"
            + "\n".join(rows)
        )


        documents.append(
            create_document(
                content,
                file_path,
                "xlsx",
                sheet_name=worksheet.title,
            )
        )


    return documents


# ============================================================
# CSV
# ============================================================

def load_csv(file_path: Path):
    """
    Carrega CSV e transforma cada arquivo em um documento textual.
    """

    documents = []


    with open(
        file_path,
        "r",
        encoding="utf-8-sig",
        newline="",
    ) as csv_file:

        reader = csv.reader(
            csv_file
        )

        rows = list(
            reader
        )


    if not rows:
        return []


    content = "\n".join(
        " | ".join(row)
        for row in rows
    )


    documents.append(
        create_document(
            content,
            file_path,
            "csv",
        )
    )


    return documents


# ============================================================
# JSON
# ============================================================

def load_json(file_path: Path):
    """
    Carrega JSON preservando sua estrutura.
    """

    with open(
        file_path,
        "r",
        encoding="utf-8",
    ) as json_file:

        data = json.load(
            json_file
        )


    content = json.dumps(
        data,
        ensure_ascii=False,
        indent=2,
    )


    return [
        create_document(
            content,
            file_path,
            "json",
        )
    ]


# ============================================================
# MARKDOWN
# ============================================================

def load_markdown(file_path: Path):
    """
    Carrega arquivos Markdown.
    """

    content = file_path.read_text(
        encoding="utf-8"
    )


    if not content.strip():
        return []


    return [
        create_document(
            content,
            file_path,
            "md",
        )
    ]


# ============================================================
# HTML
# ============================================================

def load_html(file_path: Path):
    """
    Extrai texto de páginas HTML.
    """

    html = file_path.read_text(
        encoding="utf-8"
    )


    soup = BeautifulSoup(
        html,
        "html.parser",
    )


    # Remove elementos que não são conteúdo útil
    for element in soup(
        [
            "script",
            "style",
            "noscript",
        ]
    ):

        element.decompose()


    content = soup.get_text(
        separator="\n",
        strip=True,
    )


    if not content.strip():
        return []


    return [
        create_document(
            content,
            file_path,
            "html",
        )
    ]


# ============================================================
# PPTX
# ============================================================

def load_pptx(file_path: Path):
    """
    Extrai texto dos slides de uma apresentação PowerPoint.
    """

    presentation = Presentation(
        file_path
    )


    documents = []


    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):

        texts = []


        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                text = shape.text.strip()


                if text:

                    texts.append(
                        text
                    )


        if not texts:
            continue


        content = "\n".join(
            texts
        )


        documents.append(
            create_document(
                content,
                file_path,
                "pptx",
                slide_number=slide_number,
            )
        )


    return documents


# ============================================================
# DISPATCHER
# ============================================================

def load_file(file_path: Path):
    """
    Escolhe automaticamente o loader de acordo com a extensão.
    """

    extension = file_path.suffix.lower()


    loaders = {
        ".pdf": load_pdf,
        ".docx": load_docx,
        ".xlsx": load_xlsx,
        ".csv": load_csv,
        ".json": load_json,
        ".md": load_markdown,
        ".html": load_html,
        ".pptx": load_pptx,
    }


    loader = loaders.get(
        extension
    )


    if loader is None:
        return []


    return loader(
        file_path
    )


# ============================================================
# LOADER PRINCIPAL
# ============================================================

def load_documents():
    """
    Percorre toda a pasta docs/ e carrega os documentos suportados.
    """

    if not DOCS_DIR.exists():

        raise FileNotFoundError(
            f"Pasta de documentos não encontrada: "
            f"{DOCS_DIR}"
        )


    documents = []


    files = sorted(
        file_path
        for file_path in DOCS_DIR.rglob("*")
        if (
            file_path.is_file()
            and file_path.suffix.lower()
            in SUPPORTED_EXTENSIONS
        )
    )


    print(
        f"Documentos encontrados: "
        f"{len(files)}"
    )


    for file_path in files:

        print(
            f"Carregando: "
            f"{file_path.relative_to(DOCS_DIR)}"
        )


        loaded_documents = load_file(
            file_path
        )


        # ----------------------------------------------------
        # Mostrar aviso caso o documento tenha conteúdo
        # potencialmente suspeito.
        # ----------------------------------------------------

        for document in loaded_documents:

            if document.metadata.get(
                "prompt_injection_risk",
                False
            ):

                print(
                    "  [ATENÇÃO] Possível conteúdo "
                    "de prompt injection detectado."
                )

                print(
                    "  Padrões encontrados:",
                    document.metadata.get(
                        "prompt_injection_matches",
                        []
                    )
                )


        documents.extend(
            loaded_documents
        )


    return documents


# ============================================================
# TESTE DO LOADER
# ============================================================

if __name__ == "__main__":

    documents = load_documents()


    print()

    print(
        "=" * 60
    )

    print(
        "RESULTADO DO CARREGAMENTO"
    )

    print(
        "=" * 60
    )


    print(
        f"Total de documentos carregados: "
        f"{len(documents)}"
    )


    total_characters = sum(
        len(document.page_content)
        for document in documents
    )


    print(
        f"Total de caracteres: "
        f"{total_characters}"
    )


    # ========================================================
    # RESUMO DE SEGURANÇA
    # ========================================================

    suspicious_documents = [
        document
        for document in documents
        if document.metadata.get(
            "prompt_injection_risk",
            False
        )
    ]


    print()

    print(
        "--- Verificação de segurança ---"
    )


    print(
        f"Documents com possível risco de prompt injection: "
        f"{len(suspicious_documents)}"
    )


    if suspicious_documents:

        for document in suspicious_documents:

            print()

            print(
                "Fonte:",
                document.metadata.get(
                    "source"
                )
            )

            print(
                "Padrões:",
                document.metadata.get(
                    "prompt_injection_matches"
                )
            )


    print()

    print(
        "--- Amostra do primeiro documento ---"
    )


    if documents:

        first_document = documents[0]


        print(
            first_document.page_content[:2000]
        )


        print()

        print(
            "--- Metadata ---"
        )


        print(
            first_document.metadata
        )